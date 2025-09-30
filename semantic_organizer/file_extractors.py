"""
File content extraction utilities for various file formats.

This module provides text extraction capabilities for different file types
including documents, PDFs, spreadsheets, and presentations.
"""

import logging
import mimetypes
import re
from pathlib import Path
from typing import Optional, Tuple, Dict, Any

import chardet

logger = logging.getLogger('semantic_organizer.file_extractors')


class FileExtractor:
    """Main file content extraction class that routes to appropriate extractors."""

    def __init__(self, max_file_size_mb: int = 100, sample_limit: int = 50000):
        """
        Initialize the FileExtractor.

        Args:
            max_file_size_mb: Maximum file size to process in MB
            sample_limit: Maximum characters to extract from large files
        """
        self.max_file_size = max_file_size_mb * 1024 * 1024  # Convert to bytes
        self.sample_limit = sample_limit

        # Initialize optional libraries
        self._init_libraries()

    def _init_libraries(self) -> None:
        """Initialize optional libraries and set availability flags."""
        # Import optional libraries
        try:
            import docx
            self._docx_available = True
        except ImportError:
            self._docx_available = False
            logger.warning("python-docx not available. .docx files will be processed by filename only.")

        try:
            import pdfplumber
            self._pdfplumber_available = True
        except ImportError:
            try:
                import PyPDF2
                self._pypdf2_available = True
                self._pdfplumber_available = False
            except ImportError:
                self._pypdf2_available = False
                self._pdfplumber_available = False
                logger.warning("Neither pdfplumber nor PyPDF2 available. .pdf files will be processed by filename only.")

        try:
            import openpyxl
            self._openpyxl_available = True
        except ImportError:
            self._openpyxl_available = False
            logger.warning("openpyxl not available. .xlsx files will be processed by filename only.")

        try:
            import pptx
            self._pptx_available = True
        except ImportError:
            self._pptx_available = False
            logger.warning("python-pptx not available. .pptx files will be processed by filename only.")

    def extract_text(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """
        Main dispatcher that routes to appropriate extractor.

        Args:
            file_path: Path to the file to extract text from

        Returns:
            Tuple of (extracted_text, metadata_dict)
        """
        try:
            # Check file size
            if file_path.stat().st_size > self.max_file_size:
                logger.warning(f"File {file_path} exceeds size limit ({self.max_file_size / 1024 / 1024:.1f} MB)")
                return self._fallback_extraction(file_path)

            file_type = self.get_file_type(file_path)
            metadata = {
                'file_type': file_type,
                'file_size': file_path.stat().st_size,
                'extraction_method': None,
                'content_available': True
            }

            if file_type == 'docx' and self._docx_available:
                text = self._extract_from_docx(file_path)
                metadata['extraction_method'] = 'python-docx'
            elif file_type == 'pdf':
                text = self._extract_from_pdf(file_path)
                metadata['extraction_method'] = 'pdfplumber' if self._pdfplumber_available else 'PyPDF2'
            elif file_type == 'xlsx' and self._openpyxl_available:
                text = self._extract_from_xlsx(file_path)
                metadata['extraction_method'] = 'openpyxl'
            elif file_type == 'pptx' and self._pptx_available:
                text = self._extract_from_pptx(file_path)
                metadata['extraction_method'] = 'python-pptx'
            elif file_type in ['txt', 'md']:
                text = self._extract_from_text(file_path)
                metadata['extraction_method'] = 'text_reader'
            else:
                text, metadata = self._fallback_extraction(file_path)

            # Apply sample limit if text is too long
            if len(text) > self.sample_limit:
                text = text[:self.sample_limit]
                metadata['truncated'] = True
                logger.debug(f"Truncated text from {file_path} to {self.sample_limit} characters")

            # Combine filename with content for better semantic analysis
            filename_text = file_path.stem.replace('_', ' ').replace('-', ' ')
            combined_text = f"{filename_text}. {text}"

            return combined_text, metadata

        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {e}")
            return self._fallback_extraction(file_path)

    def _extract_from_docx(self, file_path: Path) -> str:
        """Extract text from Word documents."""
        import docx

        try:
            doc = docx.Document(file_path)
            text_parts = []

            # Extract paragraph text
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text.strip())

            # Extract table text
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_parts.append(' '.join(row_text))

            return '\n'.join(text_parts)

        except Exception as e:
            logger.error(f"Error reading .docx file {file_path}: {e}")
            raise

    def _extract_from_pdf(self, file_path: Path) -> str:
        """Extract text from PDF files."""
        text_parts = []

        if self._pdfplumber_available:
            import pdfplumber

            try:
                with pdfplumber.open(file_path) as pdf:
                    max_pages = min(50, len(pdf.pages))  # Limit to first 50 pages

                    for i, page in enumerate(pdf.pages[:max_pages]):
                        page_text = page.extract_text()
                        if page_text and page_text.strip():
                            text_parts.append(page_text.strip())

                        # Stop if we've extracted enough text
                        if len('\n'.join(text_parts)) > self.sample_limit:
                            break

                return '\n'.join(text_parts)

            except Exception as e:
                logger.error(f"Error reading PDF with pdfplumber {file_path}: {e}")
                if self._pypdf2_available:
                    return self._extract_pdf_pypdf2(file_path)
                raise

        elif self._pypdf2_available:
            return self._extract_pdf_pypdf2(file_path)

        else:
            raise ImportError("No PDF extraction library available")

    def _extract_pdf_pypdf2(self, file_path: Path) -> str:
        """Extract text from PDF using PyPDF2 as fallback."""
        import PyPDF2

        text_parts = []
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                max_pages = min(50, len(pdf_reader.pages))

                for i in range(max_pages):
                    try:
                        page = pdf_reader.pages[i]
                        page_text = page.extract_text()
                        if page_text and page_text.strip():
                            text_parts.append(page_text.strip())

                        if len('\n'.join(text_parts)) > self.sample_limit:
                            break
                    except Exception as e:
                        logger.warning(f"Error extracting page {i} from {file_path}: {e}")
                        continue

            return '\n'.join(text_parts)

        except Exception as e:
            logger.error(f"Error reading PDF with PyPDF2 {file_path}: {e}")
            raise

    def _extract_from_xlsx(self, file_path: Path) -> str:
        """Extract text from Excel files."""
        import openpyxl

        try:
            workbook = openpyxl.load_workbook(file_path, read_only=True)
            text_parts = []

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_parts.append(f"Sheet: {sheet_name}")

                # Sample first 100 rows to avoid processing huge spreadsheets
                max_row = min(100, sheet.max_row) if sheet.max_row else 100

                for row in sheet.iter_rows(max_row=max_row, values_only=True):
                    row_text = []
                    for cell in row:
                        if cell is not None and str(cell).strip():
                            row_text.append(str(cell).strip())
                    if row_text:
                        text_parts.append(' '.join(row_text))

                # Stop if we've extracted enough text
                if len('\n'.join(text_parts)) > self.sample_limit:
                    break

            workbook.close()
            return '\n'.join(text_parts)

        except Exception as e:
            logger.error(f"Error reading .xlsx file {file_path}: {e}")
            raise

    def _extract_from_pptx(self, file_path: Path) -> str:
        """Extract text from PowerPoint presentations."""
        import pptx

        try:
            presentation = pptx.Presentation(file_path)
            text_parts = []

            for i, slide in enumerate(presentation.slides):
                slide_texts = []

                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())

                if slide_texts:
                    text_parts.append(f"Slide {i + 1}: " + ' '.join(slide_texts))

                # Stop if we've extracted enough text
                if len('\n'.join(text_parts)) > self.sample_limit:
                    break

            return '\n'.join(text_parts)

        except Exception as e:
            logger.error(f"Error reading .pptx file {file_path}: {e}")
            raise

    def _extract_from_text(self, file_path: Path) -> str:
        """Extract text from plain text and markdown files."""
        try:
            # Try to detect encoding
            with open(file_path, 'rb') as file:
                raw_data = file.read()
                encoding_result = chardet.detect(raw_data)
                encoding = encoding_result['encoding'] or 'utf-8'

            # Read with detected encoding
            with open(file_path, 'r', encoding=encoding) as file:
                text = file.read()

                # For markdown files, optionally clean up some formatting
                if file_path.suffix.lower() == '.md':
                    # Remove some markdown syntax but keep content
                    text = re.sub(r'^#+\s*', '', text, flags=re.MULTILINE)  # Remove headers
                    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)  # Remove bold
                    text = re.sub(r'\*(.*?)\*', r'\1', text)  # Remove italic

                return text

        except Exception as e:
            logger.error(f"Error reading text file {file_path}: {e}")
            raise

    def _fallback_extraction(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Fallback extraction using filename and basic metadata."""
        filename_text = file_path.stem.replace('_', ' ').replace('-', ' ')

        metadata = {
            'file_type': self.get_file_type(file_path),
            'file_size': file_path.stat().st_size if file_path.exists() else 0,
            'extraction_method': 'filename_only',
            'content_available': False
        }

        return filename_text, metadata

    def get_file_type(self, file_path: Path) -> str:
        """
        Determine file type from extension.

        Args:
            file_path: Path to the file

        Returns:
            File type string
        """
        suffix = file_path.suffix.lower()

        type_mapping = {
            '.txt': 'txt',
            '.md': 'md',
            '.docx': 'docx',
            '.doc': 'docx',  # Treat .doc as .docx for now
            '.pdf': 'pdf',
            '.xlsx': 'xlsx',
            '.xls': 'xlsx',  # Treat .xls as .xlsx for now
            '.pptx': 'pptx',
            '.ppt': 'pptx',  # Treat .ppt as .pptx for now
        }

        return type_mapping.get(suffix, 'unknown')

    def handle_extraction_error(self, file_path: Path, error: Exception) -> Tuple[str, Dict[str, Any]]:
        """
        Handle extraction errors and provide fallback.

        Args:
            file_path: Path to the file that failed
            error: The exception that occurred

        Returns:
            Tuple of (fallback_text, error_metadata)
        """
        logger.warning(f"Content extraction failed for {file_path}: {error}")
        return self._fallback_extraction(file_path)