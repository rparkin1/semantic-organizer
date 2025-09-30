"""
Script to create example PowerPoint, Excel, and PDF files for testing.

This script creates real office files that can be processed by the semantic organizer.
"""

import sys
from pathlib import Path


def create_excel_files(base_dir: Path):
    """Create example Excel files."""
    try:
        import openpyxl
        from openpyxl import Workbook

        # Sales Report
        wb1 = Workbook()
        ws1 = wb1.active
        ws1.title = "Sales Data"

        # Headers
        ws1['A1'] = "Product"
        ws1['B1'] = "Q1 Sales"
        ws1['C1'] = "Q2 Sales"
        ws1['D1'] = "Q3 Sales"
        ws1['E1'] = "Q4 Sales"

        # Sales data
        sales_data = [
            ["Laptops", 15000, 18000, 22000, 25000],
            ["Smartphones", 32000, 35000, 38000, 42000],
            ["Tablets", 8000, 9500, 11000, 12500],
            ["Headphones", 5000, 6000, 7500, 8500],
            ["Smart Watches", 12000, 15000, 18000, 20000]
        ]

        for i, row in enumerate(sales_data, 2):
            for j, value in enumerate(row, 1):
                ws1.cell(row=i, column=j, value=value)

        wb1.save(base_dir / "quarterly_sales_report.xlsx")

        # Budget Analysis
        wb2 = Workbook()
        ws2 = wb2.active
        ws2.title = "Budget Analysis"

        ws2['A1'] = "Department"
        ws2['B1'] = "Allocated Budget"
        ws2['C1'] = "Actual Spending"
        ws2['D1'] = "Variance"

        budget_data = [
            ["Marketing", 50000, 48500, 1500],
            ["Research & Development", 75000, 82000, -7000],
            ["Operations", 120000, 115000, 5000],
            ["Human Resources", 45000, 43000, 2000],
            ["IT Infrastructure", 60000, 58500, 1500]
        ]

        for i, row in enumerate(budget_data, 2):
            for j, value in enumerate(row, 1):
                ws2.cell(row=i, column=j, value=value)

        wb2.save(base_dir / "annual_budget_analysis.xlsx")

        # Employee Database
        wb3 = Workbook()
        ws3 = wb3.active
        ws3.title = "Employee Records"

        ws3['A1'] = "Employee ID"
        ws3['B1'] = "Name"
        ws3['C1'] = "Department"
        ws3['D1'] = "Position"
        ws3['E1'] = "Salary"
        ws3['F1'] = "Start Date"

        employee_data = [
            ["E001", "Alice Johnson", "Engineering", "Senior Developer", 95000, "2020-03-15"],
            ["E002", "Bob Smith", "Marketing", "Marketing Manager", 75000, "2019-07-01"],
            ["E003", "Carol Davis", "HR", "HR Specialist", 55000, "2021-01-10"],
            ["E004", "David Wilson", "Engineering", "DevOps Engineer", 85000, "2020-11-20"],
            ["E005", "Emma Brown", "Finance", "Financial Analyst", 65000, "2021-06-01"]
        ]

        for i, row in enumerate(employee_data, 2):
            for j, value in enumerate(row, 1):
                ws3.cell(row=i, column=j, value=value)

        wb3.save(base_dir / "employee_database.xlsx")

        print("✓ Created 3 Excel files")
        return True

    except ImportError:
        print("⚠ openpyxl not available - skipping Excel file creation")
        return False


def create_powerpoint_files(base_dir: Path):
    """Create example PowerPoint files."""
    try:
        from pptx import Presentation
        from pptx.util import Inches

        # Marketing Presentation
        prs1 = Presentation()

        # Title slide
        title_slide_layout = prs1.slide_layouts[0]
        slide = prs1.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Digital Marketing Strategy 2024"
        subtitle.text = "Driving Growth Through Innovation"

        # Content slide
        bullet_slide_layout = prs1.slide_layouts[1]
        slide = prs1.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = 'Key Strategies'
        tf = body_shape.text_frame
        tf.text = 'Social Media Marketing'

        p = tf.add_paragraph()
        p.text = 'Search Engine Optimization'
        p = tf.add_paragraph()
        p.text = 'Content Marketing'
        p = tf.add_paragraph()
        p.text = 'Email Campaigns'
        p = tf.add_paragraph()
        p.text = 'Influencer Partnerships'

        # Results slide
        slide = prs1.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = 'Expected Results'
        tf = body_shape.text_frame
        tf.text = '25% increase in website traffic'

        p = tf.add_paragraph()
        p.text = '15% growth in social media engagement'
        p = tf.add_paragraph()
        p.text = '20% improvement in conversion rate'
        p = tf.add_paragraph()
        p.text = '30% boost in brand awareness'

        prs1.save(base_dir / "marketing_strategy_presentation.pptx")

        # Technical Presentation
        prs2 = Presentation()

        # Title slide
        slide = prs2.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Machine Learning Implementation"
        subtitle.text = "AI Solutions for Business Intelligence"

        # Architecture slide
        slide = prs2.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = 'System Architecture'
        tf = body_shape.text_frame
        tf.text = 'Data Collection Layer'

        p = tf.add_paragraph()
        p.text = 'Data Processing Pipeline'
        p = tf.add_paragraph()
        p.text = 'Machine Learning Models'
        p = tf.add_paragraph()
        p.text = 'API Service Layer'
        p = tf.add_paragraph()
        p.text = 'Frontend Dashboard'

        # Technologies slide
        slide = prs2.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = 'Technology Stack'
        tf = body_shape.text_frame
        tf.text = 'Python & Scikit-learn for ML models'

        p = tf.add_paragraph()
        p.text = 'PostgreSQL for data storage'
        p = tf.add_paragraph()
        p.text = 'FastAPI for backend services'
        p = tf.add_paragraph()
        p.text = 'React for frontend interface'
        p = tf.add_paragraph()
        p.text = 'Docker for containerization'

        prs2.save(base_dir / "ml_implementation_slides.pptx")

        # Training Presentation
        prs3 = Presentation()

        # Title slide
        slide = prs3.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Employee Safety Training"
        subtitle.text = "Workplace Safety Protocols and Procedures"

        # Safety guidelines
        slide = prs3.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = 'Safety Guidelines'
        tf = body_shape.text_frame
        tf.text = 'Always wear appropriate PPE'

        p = tf.add_paragraph()
        p.text = 'Report hazards immediately'
        p = tf.add_paragraph()
        p.text = 'Follow emergency procedures'
        p = tf.add_paragraph()
        p.text = 'Keep work areas clean'
        p = tf.add_paragraph()
        p.text = 'Use equipment properly'

        prs3.save(base_dir / "safety_training_presentation.pptx")

        print("✓ Created 3 PowerPoint files")
        return True

    except ImportError:
        print("⚠ python-pptx not available - skipping PowerPoint file creation")
        return False


def create_pdf_files(base_dir: Path):
    """Create example PDF files using reportlab if available."""
    try:
        from reportlab.lib.pagesizes import letter, A4
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.lib import colors

        styles = getSampleStyleSheet()

        # Research Paper PDF
        doc1 = SimpleDocTemplate(base_dir / "ai_research_paper.pdf", pagesize=letter)
        story1 = []

        # Title
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=16,
            spaceAfter=30,
            alignment=1  # Center
        )

        story1.append(Paragraph("Advances in Natural Language Processing", title_style))
        story1.append(Spacer(1, 12))

        # Abstract
        story1.append(Paragraph("Abstract", styles['Heading2']))
        abstract_text = """
        This paper presents recent advances in natural language processing (NLP) using
        transformer-based architectures. We explore the applications of large language models
        in text classification, sentiment analysis, and machine translation. Our experiments
        show significant improvements over traditional approaches, with accuracy increases
        of up to 15% across multiple benchmarks.
        """
        story1.append(Paragraph(abstract_text, styles['Normal']))
        story1.append(Spacer(1, 12))

        # Introduction
        story1.append(Paragraph("1. Introduction", styles['Heading2']))
        intro_text = """
        Natural Language Processing has undergone a revolution with the introduction of
        transformer architectures. Models like BERT, GPT, and T5 have achieved
        state-of-the-art performance on numerous NLP tasks. This research investigates
        the practical applications of these models in real-world scenarios.
        """
        story1.append(Paragraph(intro_text, styles['Normal']))
        story1.append(Spacer(1, 12))

        # Methodology
        story1.append(Paragraph("2. Methodology", styles['Heading2']))
        method_text = """
        We conducted experiments using three different transformer models:
        BERT for classification tasks, GPT-3 for text generation, and T5 for
        translation tasks. Each model was fine-tuned on domain-specific datasets
        and evaluated using standard metrics.
        """
        story1.append(Paragraph(method_text, styles['Normal']))

        doc1.build(story1)

        # Financial Report PDF
        doc2 = SimpleDocTemplate(base_dir / "financial_quarterly_report.pdf", pagesize=A4)
        story2 = []

        story2.append(Paragraph("Q4 2024 Financial Report", title_style))
        story2.append(Spacer(1, 20))

        story2.append(Paragraph("Executive Summary", styles['Heading2']))
        exec_summary = """
        The fourth quarter of 2024 showed strong financial performance with revenue
        growth of 12% compared to Q4 2023. Net income increased by 8%, driven by
        improved operational efficiency and cost management initiatives. Key highlights
        include expansion into new markets and successful product launches.
        """
        story2.append(Paragraph(exec_summary, styles['Normal']))
        story2.append(Spacer(1, 12))

        # Financial table
        story2.append(Paragraph("Financial Highlights", styles['Heading2']))

        data = [
            ['Metric', 'Q4 2024', 'Q4 2023', 'Change'],
            ['Revenue ($M)', '125.6', '112.1', '+12.0%'],
            ['Net Income ($M)', '18.4', '17.0', '+8.2%'],
            ['Gross Margin (%)', '42.3', '40.1', '+2.2pp'],
            ['Operating Margin (%)', '15.2', '14.8', '+0.4pp']
        ]

        table = Table(data)
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 14),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))

        story2.append(table)
        story2.append(Spacer(1, 12))

        outlook_text = """
        Looking ahead to 2025, we expect continued growth driven by digital transformation
        initiatives and expansion into emerging markets. We project revenue growth of
        10-15% and maintain our commitment to innovation and operational excellence.
        """
        story2.append(Paragraph("2025 Outlook", styles['Heading2']))
        story2.append(Paragraph(outlook_text, styles['Normal']))

        doc2.build(story2)

        # User Manual PDF
        doc3 = SimpleDocTemplate(base_dir / "software_user_manual.pdf", pagesize=letter)
        story3 = []

        story3.append(Paragraph("Software User Manual", title_style))
        story3.append(Paragraph("Version 2.1", styles['Normal']))
        story3.append(Spacer(1, 20))

        story3.append(Paragraph("Table of Contents", styles['Heading2']))
        toc_text = """
        1. Installation Instructions<br/>
        2. Getting Started<br/>
        3. User Interface Overview<br/>
        4. Features and Functions<br/>
        5. Troubleshooting<br/>
        6. Technical Support
        """
        story3.append(Paragraph(toc_text, styles['Normal']))
        story3.append(Spacer(1, 12))

        story3.append(Paragraph("1. Installation Instructions", styles['Heading2']))
        install_text = """
        To install the software, please follow these steps:<br/>
        1. Download the installer from our website<br/>
        2. Run the installer as administrator<br/>
        3. Follow the setup wizard instructions<br/>
        4. Restart your computer when prompted<br/>
        5. Launch the application from the desktop shortcut
        """
        story3.append(Paragraph(install_text, styles['Normal']))
        story3.append(Spacer(1, 12))

        story3.append(Paragraph("2. Getting Started", styles['Heading2']))
        start_text = """
        After installation, the first time you launch the software, you will be
        guided through an initial setup process. This includes creating your user
        profile, configuring preferences, and connecting to your data sources.
        The setup wizard will walk you through each step.
        """
        story3.append(Paragraph(start_text, styles['Normal']))

        doc3.build(story3)

        print("✓ Created 3 PDF files")
        return True

    except ImportError:
        print("⚠ reportlab not available - skipping PDF file creation")
        return False


def create_mock_office_files(base_dir: Path):
    """Create mock office files with content that our extractors can still process."""

    # Mock PDF content (text that would be extracted)
    (base_dir / "mock_research_paper.pdf").write_text("""
    MOCK PDF CONTENT - Machine Learning in Healthcare Applications

    Abstract: This research paper explores the application of machine learning
    algorithms in healthcare diagnosis and treatment optimization. We present
    findings from clinical trials using neural networks for medical image analysis.

    Keywords: machine learning, healthcare, neural networks, medical diagnosis
    """)

    (base_dir / "mock_financial_report.pdf").write_text("""
    MOCK PDF CONTENT - Annual Financial Summary 2024

    Revenue Analysis:
    - Q1: $2.5M (+8% YoY)
    - Q2: $2.8M (+12% YoY)
    - Q3: $3.1M (+15% YoY)
    - Q4: $3.4M (+18% YoY)

    Total Annual Revenue: $11.8M
    Net Profit Margin: 22.3%
    """)

    # Mock Excel content
    (base_dir / "mock_sales_data.xlsx").write_text("""
    MOCK EXCEL CONTENT - Sales Performance Dashboard

    Product Categories:
    Electronics: $450K revenue, 1,250 units sold
    Clothing: $280K revenue, 2,100 units sold
    Home & Garden: $320K revenue, 890 units sold
    Sports Equipment: $190K revenue, 650 units sold

    Top performing region: Northeast (35% of total sales)
    Growth rate: +14% compared to previous year
    """)

    (base_dir / "mock_inventory_report.xlsx").write_text("""
    MOCK EXCEL CONTENT - Inventory Management Report

    Warehouse Status:
    - Total SKUs: 2,847
    - In Stock: 2,456 (86.3%)
    - Low Stock: 287 (10.1%)
    - Out of Stock: 104 (3.6%)

    Reorder recommendations: 391 items
    Fast-moving items: Electronics, Personal Care
    """)

    # Mock PowerPoint content
    (base_dir / "mock_business_plan.pptx").write_text("""
    MOCK POWERPOINT CONTENT - Business Expansion Strategy

    Slide 1: Executive Summary
    - Market opportunity: $2.5B addressable market
    - Competitive advantage: AI-powered analytics
    - Financial projections: 3x revenue growth in 24 months

    Slide 2: Market Analysis
    - Target demographics: Tech-savvy professionals 25-45
    - Geographic expansion: West Coast, Texas, Florida
    - Market penetration strategy: Digital marketing + partnerships

    Slide 3: Implementation Timeline
    - Phase 1 (Q1): Team expansion, technology platform
    - Phase 2 (Q2-Q3): Market entry, customer acquisition
    - Phase 3 (Q4): Scale operations, optimize performance
    """)

    (base_dir / "mock_training_slides.pptx").write_text("""
    MOCK POWERPOINT CONTENT - Employee Development Program

    Slide 1: Training Objectives
    - Skill enhancement in digital tools
    - Leadership development pathways
    - Cross-functional collaboration

    Slide 2: Program Structure
    - 8-week intensive training
    - Hands-on workshops and simulations
    - Mentorship and peer learning
    - Performance assessments

    Slide 3: Expected Outcomes
    - 25% increase in productivity metrics
    - Improved employee satisfaction scores
    - Enhanced career advancement opportunities
    """)

    print("✓ Created 6 mock office files (for systems without required libraries)")


def main():
    """Create all office files for testing."""
    base_dir = Path("test_input")
    base_dir.mkdir(exist_ok=True)

    print("Creating office files for testing...")

    created_real = []

    # Try to create real files first
    if create_excel_files(base_dir):
        created_real.append("Excel")

    if create_powerpoint_files(base_dir):
        created_real.append("PowerPoint")

    if create_pdf_files(base_dir):
        created_real.append("PDF")

    # Always create mock files as backup/additional content
    create_mock_office_files(base_dir)

    print(f"\n✅ Office file creation complete!")
    if created_real:
        print(f"Real files created: {', '.join(created_real)}")
    print("Mock files created for all formats")

    print(f"\nTotal files in test_input: {len(list(base_dir.glob('*')))}")

    print("\nTo test with these files, run:")
    print("python -m semantic_organizer test_input test_output --dry-run --verbose")


if __name__ == "__main__":
    main()