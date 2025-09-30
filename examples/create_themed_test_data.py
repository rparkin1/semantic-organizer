"""
Create focused test data with clear semantic themes for better clustering.

Each theme will have 3-5 files with similar content to test semantic similarity.
"""

from pathlib import Path
import shutil


def create_machine_learning_theme(base_dir: Path):
    """Create files focused on machine learning and AI."""

    # Research papers
    (base_dir / "deep_learning_neural_networks.txt").write_text("""
    Deep Learning and Neural Network Architectures

    This paper explores various deep learning architectures including convolutional neural networks (CNNs),
    recurrent neural networks (RNNs), and transformer models. We discuss backpropagation algorithms,
    gradient descent optimization, and regularization techniques. Applications include computer vision,
    natural language processing, and speech recognition.

    Keywords: deep learning, neural networks, CNN, RNN, transformers, backpropagation
    """)

    (base_dir / "machine_learning_algorithms_comparison.txt").write_text("""
    Comparison of Machine Learning Algorithms

    A comprehensive analysis of supervised and unsupervised learning algorithms. We compare decision trees,
    random forests, support vector machines, k-means clustering, and neural networks. Performance metrics
    include accuracy, precision, recall, F1-score, and computational complexity.

    Topics covered: supervised learning, unsupervised learning, classification, regression, clustering
    """)

    (base_dir / "ai_computer_vision_applications.txt").write_text("""
    Artificial Intelligence in Computer Vision

    Computer vision applications using AI and deep learning techniques. Object detection using YOLO,
    image classification with ResNet, semantic segmentation with U-Net, and facial recognition systems.
    Implementation details include data preprocessing, model training, and performance evaluation.

    Applications: object detection, image classification, face recognition, autonomous vehicles
    """)

    (base_dir / "natural_language_processing_transformers.txt").write_text("""
    Natural Language Processing with Transformer Models

    Advanced NLP techniques using BERT, GPT, and T5 transformer architectures. Text classification,
    sentiment analysis, named entity recognition, and machine translation. Pre-training and fine-tuning
    strategies for domain-specific applications.

    Technologies: BERT, GPT, T5, attention mechanisms, transfer learning
    """)


def create_financial_theme(base_dir: Path):
    """Create files focused on finance and business."""

    (base_dir / "quarterly_financial_report_q1.txt").write_text("""
    Q1 2024 Financial Performance Report

    Revenue: $12.5M (+15% YoY)
    Net Income: $2.8M (+18% YoY)
    EBITDA: $3.2M
    Operating Margin: 22.4%

    Key Highlights:
    - Strong growth in subscription revenue
    - Expansion into European markets
    - Cost optimization initiatives successful
    - Cash flow positive for 6th consecutive quarter

    Market Analysis: Technology sector showing resilience despite economic headwinds
    """)

    (base_dir / "investment_portfolio_analysis_2024.txt").write_text("""
    Investment Portfolio Performance Analysis 2024

    Portfolio Value: $2.4M (+12% YTD)
    Asset Allocation:
    - Equities: 65% (Technology, Healthcare, Financial Services)
    - Fixed Income: 25% (Government bonds, Corporate bonds)
    - Alternative Investments: 10% (REITs, Commodities)

    Performance Metrics:
    - Sharpe Ratio: 1.34
    - Maximum Drawdown: -8.2%
    - Beta: 0.89

    Recommendation: Maintain current allocation, consider increasing exposure to emerging markets
    """)

    (base_dir / "budget_forecast_annual_planning.txt").write_text("""
    Annual Budget Forecast and Strategic Planning

    2024 Budget Summary:
    Total Revenue Target: $48M
    Operating Expenses: $36M
    Capital Expenditure: $4M
    Net Profit Target: $8M

    Department Budgets:
    - Sales & Marketing: $12M
    - Research & Development: $8M
    - Operations: $10M
    - Administration: $6M

    Growth Initiatives: Product development, market expansion, team scaling
    """)

    (base_dir / "financial_risk_assessment_report.txt").write_text("""
    Financial Risk Assessment and Management

    Risk Categories:
    - Market Risk: Currency fluctuations, interest rate changes
    - Credit Risk: Customer default probability assessment
    - Operational Risk: Process failures, system downtime
    - Liquidity Risk: Cash flow management analysis

    Mitigation Strategies:
    - Diversified revenue streams
    - Hedging instruments for currency exposure
    - Robust credit scoring models
    - Emergency fund maintenance (6 months operating expenses)

    Risk Score: Medium (2.3/5.0)
    """)


def create_recipe_cooking_theme(base_dir: Path):
    """Create files focused on cooking and recipes."""

    (base_dir / "italian_pasta_recipes_collection.txt").write_text("""
    Traditional Italian Pasta Recipes

    1. Spaghetti Carbonara
    Ingredients: spaghetti, eggs, pecorino romano, pancetta, black pepper
    Technique: Create silky egg sauce without scrambling

    2. Fettuccine Alfredo
    Ingredients: fettuccine, butter, parmigiano-reggiano, heavy cream
    Method: Toss hot pasta with butter and cheese for creamy coating

    3. Penne Arrabbiata
    Ingredients: penne, tomatoes, garlic, red chili flakes, olive oil
    Preparation: Spicy tomato sauce with fresh herbs

    Cooking Tips: Use pasta water for sauce consistency, always cook al dente
    """)

    (base_dir / "baking_bread_sourdough_techniques.txt").write_text("""
    Artisanal Bread Baking and Sourdough Techniques

    Sourdough Starter Maintenance:
    - Daily feeding with equal parts flour and water
    - Maintain at room temperature or refrigerate
    - Look for doubling in size and pleasant tangy aroma

    Bread Making Process:
    1. Autolyse: Mix flour and water, let rest 30 minutes
    2. Add starter and salt, perform stretch and folds
    3. Bulk fermentation: 4-6 hours with periodic folding
    4. Shape and final proof: 2-4 hours or overnight
    5. Bake in Dutch oven at 450°F with steam

    Troubleshooting: Dense bread, over-proofing, temperature control
    """)

    (base_dir / "healthy_meal_prep_recipes.txt").write_text("""
    Weekly Healthy Meal Prep Recipes

    High-Protein Options:
    - Grilled chicken with quinoa and roasted vegetables
    - Baked salmon with sweet potato and broccoli
    - Turkey meatballs with zucchini noodles
    - Tofu stir-fry with brown rice and mixed vegetables

    Meal Prep Tips:
    - Cook proteins in bulk on Sunday
    - Pre-cut vegetables for easy assembly
    - Use glass containers for storage
    - Label with dates and contents

    Nutritional Balance: 40% carbs, 30% protein, 30% healthy fats
    Storage: 3-4 days refrigerated, up to 3 months frozen
    """)

    (base_dir / "international_cuisine_fusion_dishes.txt").write_text("""
    International Cuisine and Fusion Cooking

    Asian-Inspired Dishes:
    - Korean BBQ tacos with kimchi slaw
    - Thai green curry risotto
    - Japanese ramen burger
    - Chinese five-spice roasted duck

    Mediterranean Fusion:
    - Greek tzatziki pasta salad
    - Moroccan tagine pizza
    - Spanish paella arancini
    - Turkish döner kebab wrap

    Cooking Techniques: Wok hei, slow braising, spice blooming, fermentation
    Ingredients: Miso paste, harissa, sumac, fish sauce, preserved lemons
    """)


def create_travel_theme(base_dir: Path):
    """Create files focused on travel and destinations."""

    (base_dir / "europe_backpacking_complete_guide.txt").write_text("""
    Complete Guide to Backpacking Through Europe

    Essential Countries to Visit:
    - France: Paris, Lyon, Nice - Art, cuisine, Mediterranean coast
    - Italy: Rome, Florence, Venice - History, Renaissance art, canals
    - Germany: Berlin, Munich, Hamburg - Culture, beer gardens, nightlife
    - Spain: Madrid, Barcelona, Seville - Architecture, flamenco, tapas
    - Netherlands: Amsterdam, Utrecht - Canals, museums, cycling culture

    Budget Planning:
    - Accommodation: €20-40/night (hostels)
    - Transportation: Eurail pass €400-600
    - Food: €25-35/day
    - Activities: €10-20/day per attraction

    Packing Essentials: Comfortable walking shoes, lightweight clothes, portable charger
    """)

    (base_dir / "asia_cultural_immersion_travel.txt").write_text("""
    Cultural Immersion Travel in Southeast Asia

    Destinations and Experiences:
    - Thailand: Buddhist temples, floating markets, Thai cooking classes
    - Vietnam: Motorbike tours, pho cooking, Halong Bay cruising
    - Cambodia: Angkor Wat sunrise, Khmer culture, traditional crafts
    - Laos: Mekong River journey, monks' alms ceremony, textile weaving
    - Indonesia: Balinese ceremonies, volcano hiking, traditional dance

    Cultural Etiquette:
    - Remove shoes when entering homes and temples
    - Dress modestly at religious sites
    - Learn basic greetings in local languages
    - Respect for elders and monks

    Best Time to Visit: November to March (dry season)
    """)

    (base_dir / "adventure_travel_extreme_sports.txt").write_text("""
    Adventure Travel and Extreme Sports Destinations

    Mountain Adventures:
    - Nepal: Everest Base Camp trekking, mountain climbing
    - Patagonia: Glacier hiking, ice climbing, wind surfing
    - Swiss Alps: Skiing, mountaineering, paragliding
    - Canadian Rockies: Backcountry camping, wildlife photography

    Water Sports:
    - New Zealand: Bungee jumping, white water rafting, skydiving
    - Costa Rica: Zip-lining, surfing, volcano boarding
    - Australia: Great Barrier Reef diving, shark cage diving
    - Hawaii: Big wave surfing, volcano exploration

    Safety Considerations: Proper equipment, local guides, insurance coverage
    Physical Preparation: Fitness training, altitude acclimatization
    """)

    (base_dir / "luxury_resort_hotel_reviews.txt").write_text("""
    Luxury Resort and Hotel Reviews Worldwide

    Tropical Destinations:
    - Maldives: Overwater bungalows, private butler service, world-class spa
    - Bora Bora: Lagoon views, French Polynesian cuisine, romantic sunsets
    - Seychelles: Private beaches, conservation programs, Creole culture
    - Fiji: Island hopping, coral reef snorkeling, traditional ceremonies

    City Hotels:
    - Tokyo: Minimalist design, Michelin-starred restaurants, rooftop bars
    - Paris: Historic châteaux, fashion district access, wine tastings
    - New York: Central Park views, Broadway packages, fine dining
    - London: Royal connections, afternoon tea service, theater district

    Amenities: Infinity pools, private beaches, concierge services, spa treatments
    """)


def create_fitness_health_theme(base_dir: Path):
    """Create files focused on fitness and health."""

    (base_dir / "strength_training_workout_programs.txt").write_text("""
    Comprehensive Strength Training Programs

    Beginner Program (3 days/week):
    Day A: Squats 3x8-10, Bench Press 3x8-10, Bent-over Rows 3x8-10
    Day B: Deadlifts 3x5, Overhead Press 3x8-10, Pull-ups 3xAMRAP
    Day C: Squats 3x8-10, Incline Press 3x8-10, Lat Pulldowns 3x10-12

    Progressive Overload:
    - Increase weight by 2.5-5 lbs when you can complete all sets/reps
    - Track workouts in a log book or app
    - Focus on proper form before adding weight

    Nutrition: 1g protein per lb bodyweight, adequate carbohydrates for energy
    Recovery: 48 hours rest between training same muscle groups
    """)

    (base_dir / "cardio_endurance_training_guide.txt").write_text("""
    Cardiovascular Endurance Training Guide

    Training Zones:
    - Zone 1 (60-70% max HR): Easy conversational pace
    - Zone 2 (70-80% max HR): Aerobic base building
    - Zone 3 (80-90% max HR): Lactate threshold training
    - Zone 4 (90-95% max HR): VO2 max intervals
    - Zone 5 (95%+ max HR): Neuromuscular power

    Weekly Structure:
    - 80% easy/moderate intensity (Zones 1-2)
    - 20% high intensity (Zones 3-5)
    - Include one long session per week
    - One complete rest day

    Activities: Running, cycling, swimming, rowing, elliptical
    Monitoring: Heart rate monitor, perceived exertion scale
    """)

    (base_dir / "nutrition_meal_planning_macros.txt").write_text("""
    Nutrition and Macro-Based Meal Planning

    Macronutrient Ratios by Goal:
    Fat Loss: 40% protein, 35% carbs, 25% fats
    Muscle Gain: 30% protein, 45% carbs, 25% fats
    Maintenance: 25% protein, 45% carbs, 30% fats
    Athletic Performance: 20% protein, 55% carbs, 25% fats

    Meal Timing:
    - Pre-workout: Carbs + small amount protein (1-2 hours before)
    - Post-workout: Protein + carbs (within 30 minutes)
    - Spread protein intake evenly throughout the day

    Food Sources:
    Proteins: Chicken, fish, eggs, Greek yogurt, legumes
    Carbs: Oats, rice, sweet potatoes, fruits, vegetables
    Fats: Nuts, olive oil, avocado, fatty fish
    """)

    (base_dir / "mental_health_wellness_practices.txt").write_text("""
    Mental Health and Wellness Practices

    Stress Management Techniques:
    - Mindfulness meditation: 10-20 minutes daily
    - Deep breathing exercises: 4-7-8 technique
    - Progressive muscle relaxation
    - Journaling: Gratitude and reflection practices
    - Nature exposure: Minimum 20 minutes outdoors daily

    Sleep Optimization:
    - Consistent sleep schedule (7-9 hours)
    - Cool, dark, quiet environment
    - Limit screens 1 hour before bed
    - Avoid caffeine after 2 PM
    - Relaxation routine before sleep

    Social Connections: Regular contact with friends/family, community involvement
    Professional Help: Therapy, counseling, support groups when needed
    """)


def create_real_office_files(base_dir: Path):
    """Create realistic office documents using available libraries."""

    # Try to create real Excel files
    try:
        import openpyxl
        from openpyxl import Workbook
        from openpyxl.styles import Font, Alignment, PatternFill

        # ML Research Data
        wb_ml = Workbook()
        ws_ml = wb_ml.active
        ws_ml.title = "ML Model Performance"

        headers = ["Model", "Accuracy", "Precision", "Recall", "F1-Score", "Training Time"]
        for col, header in enumerate(headers, 1):
            cell = ws_ml.cell(row=1, column=col, value=header)
            cell.font = Font(bold=True)
            cell.fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")

        ml_data = [
            ["Random Forest", 0.87, 0.85, 0.89, 0.87, "2.3 min"],
            ["SVM", 0.83, 0.81, 0.86, 0.83, "5.7 min"],
            ["Neural Network", 0.91, 0.89, 0.93, 0.91, "15.2 min"],
            ["Logistic Regression", 0.79, 0.77, 0.82, 0.79, "0.8 min"],
            ["XGBoost", 0.89, 0.87, 0.91, 0.89, "3.1 min"]
        ]

        for row_idx, row_data in enumerate(ml_data, 2):
            for col_idx, value in enumerate(row_data, 1):
                ws_ml.cell(row=row_idx, column=col_idx, value=value)

        wb_ml.save(base_dir / "machine_learning_model_comparison.xlsx")

        # Financial Budget
        wb_fin = Workbook()
        ws_fin = wb_fin.active
        ws_fin.title = "Q1 Financial Results"

        fin_headers = ["Department", "Budget", "Actual", "Variance", "Percentage"]
        for col, header in enumerate(fin_headers, 1):
            cell = ws_fin.cell(row=1, column=col, value=header)
            cell.font = Font(bold=True)
            cell.fill = PatternFill(start_color="4F81BD", end_color="4F81BD", fill_type="solid")

        fin_data = [
            ["Marketing", 50000, 47500, 2500, "95%"],
            ["R&D", 75000, 78000, -3000, "104%"],
            ["Sales", 40000, 38500, 1500, "96%"],
            ["Operations", 60000, 59200, 800, "99%"],
            ["IT", 30000, 31200, -1200, "104%"]
        ]

        for row_idx, row_data in enumerate(fin_data, 2):
            for col_idx, value in enumerate(row_data, 1):
                ws_fin.cell(row=row_idx, column=col_idx, value=value)

        wb_fin.save(base_dir / "quarterly_financial_budget_analysis.xlsx")

        print("✅ Created real Excel files")

    except ImportError:
        print("⚠ openpyxl not available, skipping Excel files")

    # Try to create real PowerPoint files
    try:
        from pptx import Presentation
        from pptx.util import Inches

        # ML Presentation
        prs_ml = Presentation()

        # Title slide
        title_slide = prs_ml.slides.add_slide(prs_ml.slide_layouts[0])
        title_slide.shapes.title.text = "Machine Learning Project Results"
        title_slide.placeholders[1].text = "Deep Learning Model Performance Analysis"

        # Content slide
        content_slide = prs_ml.slides.add_slide(prs_ml.slide_layouts[1])
        content_slide.shapes.title.text = "Model Performance Comparison"
        content_slide.placeholders[1].text = """• Neural Networks achieved 91% accuracy
• Random Forest showed good balance of speed and performance
• XGBoost provided excellent feature importance insights
• SVM struggled with large dataset size
• Future work: ensemble methods and hyperparameter tuning"""

        prs_ml.save(base_dir / "ai_machine_learning_results_presentation.pptx")

        # Travel Presentation
        prs_travel = Presentation()

        title_slide = prs_travel.slides.add_slide(prs_travel.slide_layouts[0])
        title_slide.shapes.title.text = "European Backpacking Adventure"
        title_slide.placeholders[1].text = "3-Month Journey Through 12 Countries"

        content_slide = prs_travel.slides.add_slide(prs_travel.slide_layouts[1])
        content_slide.shapes.title.text = "Highlights and Recommendations"
        content_slide.placeholders[1].text = """• Best experiences: Northern Lights in Iceland, Venice canals
• Budget breakdown: €85/day average across all countries
• Transportation: Eurail pass worth the investment
• Accommodation: Mix of hostels and Airbnb worked well
• Cultural insights: Each country's unique personality"""

        prs_travel.save(base_dir / "european_travel_adventure_slideshow.pptx")

        print("✅ Created real PowerPoint files")

    except ImportError:
        print("⚠ python-pptx not available, skipping PowerPoint files")

    # Try to create real PDF files
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet

        styles = getSampleStyleSheet()

        # Fitness PDF
        doc_fitness = SimpleDocTemplate(str(base_dir / "fitness_training_comprehensive_guide.pdf"), pagesize=letter)
        story_fitness = []

        story_fitness.append(Paragraph("Complete Fitness Training Guide", styles['Title']))
        story_fitness.append(Spacer(1, 20))

        story_fitness.append(Paragraph("Strength Training Fundamentals", styles['Heading1']))
        fitness_content = """
        Effective strength training requires progressive overload, compound movements, and adequate recovery.
        Focus on major movement patterns: squat, deadlift, bench press, and overhead press. These exercises
        work multiple muscle groups simultaneously and provide the best return on investment for muscle growth
        and strength development. Track your workouts and gradually increase weight, reps, or sets over time.
        """
        story_fitness.append(Paragraph(fitness_content, styles['Normal']))

        story_fitness.append(Spacer(1, 12))
        story_fitness.append(Paragraph("Cardiovascular Training", styles['Heading1']))
        cardio_content = """
        Cardiovascular fitness improves heart health, endurance, and aids in recovery. Include both steady-state
        cardio and high-intensity interval training (HIIT). Start with 20-30 minutes of moderate activity
        3 times per week, gradually increasing duration and intensity based on your fitness level.
        """
        story_fitness.append(Paragraph(cardio_content, styles['Normal']))

        doc_fitness.build(story_fitness)

        print("✅ Created real PDF files")

    except ImportError:
        print("⚠ reportlab not available, skipping PDF files")


def main():
    """Create comprehensive themed test data."""
    base_dir = Path("test_input_themed")

    # Clear existing directory
    if base_dir.exists():
        shutil.rmtree(base_dir)
    base_dir.mkdir()

    print("Creating themed test data with clear semantic clusters...")

    # Create themed content
    create_machine_learning_theme(base_dir)
    print("✅ Created Machine Learning & AI theme (4 files)")

    create_financial_theme(base_dir)
    print("✅ Created Financial & Business theme (4 files)")

    create_recipe_cooking_theme(base_dir)
    print("✅ Created Cooking & Recipes theme (4 files)")

    create_travel_theme(base_dir)
    print("✅ Created Travel & Adventure theme (4 files)")

    create_fitness_health_theme(base_dir)
    print("✅ Created Fitness & Health theme (4 files)")

    # Create office documents
    create_real_office_files(base_dir)

    total_files = len(list(base_dir.glob('*.txt'))) + len(list(base_dir.glob('*.xlsx'))) + len(list(base_dir.glob('*.pptx'))) + len(list(base_dir.glob('*.pdf')))

    print(f"\n🎯 Created {total_files} themed test files in {base_dir}")
    print("\n📊 Theme Distribution:")
    print("  • Machine Learning & AI: 4+ files")
    print("  • Financial & Business: 4+ files")
    print("  • Cooking & Recipes: 4+ files")
    print("  • Travel & Adventure: 4+ files")
    print("  • Fitness & Health: 4+ files")
    print("  • Office Documents: Mixed themes")

    print(f"\n🚀 Test semantic clustering with:")
    print(f"uv run python -m semantic_organizer {base_dir} test_output_themed --dry-run --verbose")
    print("uv run python -m semantic_organizer test_input_themed test_output_themed --similarity-threshold 0.5 --dry-run")


if __name__ == "__main__":
    main()